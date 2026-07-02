#!/usr/bin/env python3
"""Generate the Chapters 4 & 5 defense deck (~15 min) as a .pptx.

On-slide text + speaker notes for every slide, figures embedded.
Output: defense_ch4_5_slides.pptx in the repo root.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    from PIL import Image
    HAVE_PIL = True
except Exception:
    HAVE_PIL = False

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG = os.path.join(ROOT, "data", "figures")

# ---- palette ----
NAVY   = RGBColor(0x10, 0x2A, 0x43)   # titles / bar
TEAL   = RGBColor(0x0E, 0x7C, 0x86)   # accent
DARK   = RGBColor(0x22, 0x2A, 0x33)   # body text
GREY   = RGBColor(0x5A, 0x66, 0x72)   # subtle
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF2, 0xF5, 0xF7)   # table alt row
RED     = RGBColor(0xB4, 0x23, 0x1A)  # emphasis (failure)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def img_ratio(path):
    """height/width ratio of an image, default ~0.7."""
    if HAVE_PIL and os.path.exists(path):
        with Image.open(path) as im:
            w, h = im.size
            return h / w
    return 0.72


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def title_bar(slide, title, kicker=None):
    """Top accent bar with title; optional small kicker (section label)."""
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.55); tf.margin_right = Inches(0.5)
    tf.margin_top = Inches(0.1); tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if kicker:
        rk = p.add_run(); rk.text = kicker.upper() + "   "
        rk.font.size = Pt(13); rk.font.bold = True; rk.font.color.rgb = RGBColor(0x6FB,0,0) if False else RGBColor(0x7F,0xC9,0xD0)
    r = p.add_run(); r.text = title
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = WHITE
    # thin teal underline accent
    acc = slide.shapes.add_shape(1, 0, Inches(1.15), SW, Inches(0.06))
    acc.fill.solid(); acc.fill.fore_color.rgb = TEAL
    acc.line.fill.background(); acc.shadow.inherit = False
    return bar


def bullets(slide, items, left=Inches(0.7), top=Inches(1.5),
            width=None, height=None, size=20, gap=10):
    """items: list of (text, level, bold, color) or plain str."""
    width = width or (SW - Inches(1.4))
    height = height or (SH - top - Inches(0.5))
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, str):
            text, level, bold, color = it, 0, False, DARK
        else:
            text = it[0]
            level = it[1] if len(it) > 1 else 0
            bold = it[2] if len(it) > 2 else False
            color = it[3] if len(it) > 3 else DARK
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        # bullet glyph via prefix
        run = p.add_run()
        prefix = "" if level == 0 else ""
        bullet = "▸ " if level == 0 else "– "
        run.text = bullet + text
        run.font.size = Pt(size - (level * 2))
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_image_fit(slide, path, left, top, max_w, max_h, center=True):
    """Place image scaled to fit max_w x max_h, centered in that box."""
    r = img_ratio(path)
    w = max_w
    h = int(w * r)
    if h > max_h:
        h = max_h
        w = int(h / r)
    x = left + (max_w - w) // 2 if center else left
    y = top + (max_h - h) // 2 if center else top
    slide.shapes.add_picture(path, x, y, width=w, height=h)


def caption(slide, text, left, top, width):
    box = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(11); r.font.italic = True; r.font.color.rgb = GREY


# ============================================================== SLIDE 1 — TITLE
s = add_slide()
bg = s.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background(); bg.shadow.inherit = False
strip = s.shapes.add_shape(1, 0, Inches(4.55), SW, Inches(0.07))
strip.fill.solid(); strip.fill.fore_color.rgb = TEAL; strip.line.fill.background(); strip.shadow.inherit = False
tb = s.shapes.add_textbox(Inches(0.9), Inches(1.8), SW - Inches(1.8), Inches(2.6))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Stochastic Path Planning for IoT-Enabled Mobile Robots in Patient Emergency Response"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
sub = s.shapes.add_textbox(Inches(0.9), Inches(4.75), SW - Inches(1.8), Inches(1.8))
tf = sub.text_frame; tf.word_wrap = True
for i, (txt, sz, col, bold) in enumerate([
    ("Paula Oyinyechukwu Irabor  ·  22120612999", 22, RGBColor(0xCF,0xE3,0xE6), True),
    ("Final Year Project Defense", 18, RGBColor(0x9F,0xB6,0xBC), False),
    ("Chapters 4 & 5 — Proposed System, Evaluation & Results", 18, RGBColor(0x9F,0xB6,0xBC), False),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6)
    run = p.add_run(); run.text = txt
    run.font.size = Pt(sz); run.font.color.rgb = col; run.font.bold = bold
set_notes(s, "Good morning. My name is Paula Irabor. My project is on stochastic path planning "
            "for IoT-enabled mobile robots responding to patient emergencies. Today I'll briefly "
            "recap Chapters 1 to 3, then focus on what's new — the system I built and the results, "
            "in Chapters 4 and 5.")

# ============================================================== SLIDE 2 — OUTLINE
s = add_slide(); title_bar(s, "Outline")
bullets(s, [
    ("Recap: Problem · Aim · Objectives · Methodology", 0, True),
    ("The proposed system & its components", 0, True),
    ("Guided walkthrough of the system in action", 0, True, TEAL),
    ("Evaluation and results", 0, True),
    ("Outcomes mapped to objectives · Conclusion", 0, True),
    ("Q&A", 0, True),
], top=Inches(1.7), size=24, gap=18)
set_notes(s, "Here's the plan. I'll keep the first three chapters short, walk you through the system "
            "visually, then spend most of our time on the results and what they mean. We'll finish "
            "with a conclusion and your questions.")

# ============================================================== SLIDE 3 — PROBLEM
s = add_slide(); title_bar(s, "Problem", kicker="Recap · Ch 1–3")
bullets(s, [
    ("IoT emergency detection and robot navigation have evolved separately", 0, False),
    ("little work connects a sensor that detects an emergency to a robot that acts on it", 1, False, GREY),
    ("Most mobile robots assume a static environment", 0, False),
    ("but a hospital is dynamic and crowded — people constantly cross the robot's path", 1, False, GREY),
    ("In an emergency, a late response is a failed response", 0, True, RED),
], top=Inches(1.7), size=22, gap=14)
set_notes(s, "The problem is twofold. First, IoT health-sensing and robot navigation are usually "
            "studied separately — there's little work connecting a sensor that detects an emergency "
            "to a robot that acts on it. Second, most navigation research assumes a static world, "
            "but a hospital is the opposite: busy corridors, people always moving. And in an "
            "emergency, a response that arrives too late is the same as no response at all. That's "
            "the gap I set out to close.")

# ============================================================== SLIDE 4 — AIM
s = add_slide(); title_bar(s, "Aim", kicker="Recap · Ch 1–3")
box = s.shapes.add_textbox(Inches(1.2), Inches(2.6), SW - Inches(2.4), Inches(2.2))
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
for txt, bold, col in [
    ("To design, implement and evaluate ", True, NAVY),
    ("an IoT-enabled mobile robot that uses ", False, DARK),
    ("stochastic path planning ", True, TEAL),
    ("to respond rapidly to patient emergencies in a simulated hospital.", False, DARK),
]:
    r = p.add_run(); r.text = txt; r.font.size = Pt(30); r.font.bold = bold; r.font.color.rgb = col
set_notes(s, "So the aim was to design, build, and evaluate a complete system: an IoT-enabled mobile "
            "robot that uses stochastic path planning to reach a patient quickly, tested in a "
            "realistic simulated hospital.")

# ============================================================== SLIDE 5 — OBJECTIVES
s = add_slide(); title_bar(s, "Objectives", kicker="Recap · Ch 1–3")
bullets(s, [
    ("1.  Implement and compare two planners — A* (deterministic) vs RRT* (stochastic)", 0, True),
    ("2.  Integrate IoT emergency detection with robot navigation, end-to-end", 0, True),
    ("3.  Evaluate on: response time, path quality, replanning efficiency, success rate", 0, True),
    ("I return to these three objectives at the end and show each was met.", 0, False, GREY),
], top=Inches(1.9), size=23, gap=20)
set_notes(s, "Three objectives. One — implement and compare two path planners: A-star, the classic "
            "deterministic one, against RRT-star, a stochastic sampling planner. Two — integrate the "
            "IoT detection with the navigation so it works end to end. Three — evaluate them fairly "
            "on four metrics. I'll come back to these three objectives at the very end and show each "
            "was met.")

# ============================================================== SLIDE 6 — METHODOLOGY
s = add_slide(); title_bar(s, "Methodology", kicker="Recap · Ch 1–3")
bullets(s, [
    ("Experimental research methodology", 0, True),
    ("ROS2 Jazzy (middleware) · Gazebo Harmonic (simulation) · Python / C++", 0, False),
    ("Compared under identical conditions: 3 crowding levels (2 / 4 / 7 humans), 10 trials each", 0, False),
    ("Success judged against a 60-second clinical deadline", 0, True, TEAL),
    ("Full 180-trial matrix run headless for reproducibility", 0, False),
], top=Inches(1.7), size=22, gap=14)
set_notes(s, "The methodology is experimental. I built everything in ROS2 Jazzy with Gazebo Harmonic "
            "for the simulation. The key to a fair comparison is identical conditions: both planners "
            "face the same hospital, the same start and goal, and three levels of crowding — two, "
            "four, and seven moving people — with ten trials each. I judged success against a "
            "sixty-second clinical deadline. I ran the full 180-trial matrix headless so the results "
            "are reproducible.")

# ============================================================== SLIDE 7 — PROPOSED SYSTEM
s = add_slide(); title_bar(s, "The Proposed System", kicker="Ch 4")
bullets(s, [
    ("An integrated ROS2 + Gazebo pipeline in three layers:", 0, True),
    ("IoT Sensing — wearable detects a fall, raises an alert with the patient's location", 1, False),
    ("Integration — ROS2 middleware syncs sensor + robot state, logs every timestamp", 1, False),
    ("Navigation — plans a route, replans around humans, drives to the patient", 1, False),
], left=Inches(0.6), top=Inches(1.6), width=Inches(6.7), size=19, gap=12)
add_image_fit(s, os.path.join(FIG, "rqt_graph.png"), Inches(7.4), Inches(1.6),
              Inches(5.6), Inches(5.0))
caption(s, "Live ROS2 node graph — a running system, not a diagram", Inches(7.4), Inches(6.7), Inches(5.6))
set_notes(s, "Here's the system. It has three layers. Sensing: a simulated wearable detects a fall and "
            "broadcasts an alert carrying the patient's location. Integration: the ROS2 middleware "
            "ties the sensor and robot together and timestamps everything for the metrics. "
            "Navigation: the robot plans a route, replans when someone blocks it, and drives to the "
            "patient. The graph on the right is the actual node graph from a live run — this is a "
            "working system, not a diagram.")

# ============================================================== SLIDE 8 — COMPONENTS
s = add_slide(); title_bar(s, "Components — the Pipeline", kicker="Ch 4")
bullets(s, [
    ("emergency_sensing  →  fires on a fall  →  publishes emergency_alert (id, type, location)", 0, False),
    ("path_planner  →  runs A* or RRT* via one planner_mode switch", 0, False),
    ("robot_control  →  drives the path, replans on obstacles", 0, False),
    ("human_obstacles  →  scripted moving people cross the robot's path", 0, False),
    ("response_logger  →  timestamps everything for the metrics", 0, False),
], top=Inches(1.8), size=21, gap=18)
set_notes(s, "Breaking that into components: the sensing node fires on a fall and publishes the "
            "alert. The path-planner node runs either A-star or RRT-star — switched by a single "
            "parameter, so the comparison is truly apples-to-apples. The control node drives the "
            "robot and triggers replanning. A human-obstacles node adds the moving people. And the "
            "logger records every timestamp that the evaluation needs.")

# ============================================================== SLIDE 9 — THE TWO PLANNERS
s = add_slide(); title_bar(s, "The Two Planners", kicker="Ch 4")
bullets(s, [
    ("A* — 8-connected grid, octile heuristic  →  cheap, deterministic", 0, True),
    ("RRT* — 25% goal bias, shrinking rewire radius, 0.8 m steer  →  stochastic, explores freely", 0, True),
    ("Fair comparison: identical start / goal; all humans walk at the same speed", 0, False),
    ("crowding means more people, not faster people — isolating the planner as the only variable", 1, False, GREY),
], top=Inches(1.8), size=21, gap=16)
set_notes(s, "The two planners. A-star searches a grid with a standard heuristic — it's cheap and "
            "always returns the same optimal path. RRT-star is different: it randomly samples the "
            "space and rewires toward better routes, so it explores freely. To keep the comparison "
            "fair, both get identical start and goal positions, and crucially the humans always walk "
            "at the same speed — more crowding means more people, not faster ones. That isolates the "
            "planner as the only variable.")

# ============================================================== SLIDE 10 — ENVIRONMENT
s = add_slide(); title_bar(s, "The Hospital Environment & Scenario", kicker="Walkthrough")
# three world figures across the bottom
imgs = ["world_low_crowding.png", "world_moderate_crowding.png", "world_high_crowding.png"]
labels = ["Low — 2 humans", "Moderate — 4 humans", "High — 7 humans"]
colw = Inches(4.25); gap = Inches(0.12); startx = Inches(0.25); topy = Inches(2.55)
for i, (im, lb) in enumerate(zip(imgs, labels)):
    x = startx + i * (colw + gap)
    add_image_fit(s, os.path.join(FIG, im), x, topy, colw, Inches(3.9))
    caption(s, lb, x, Inches(6.65), colw)
bullets(s, [
    ("16 × 14 m two-room hospital · robot start (nurse station) · near patient (a) and far patient (b)", 0, False),
    ("Crowding rises left → right; the far-patient route crosses the busiest corridor — where planners diverge", 0, True, TEAL),
], top=Inches(1.55), size=17, gap=8)
set_notes(s, "Now a guided walkthrough of the system in action. This is the hospital — a "
            "sixteen-by-fourteen-metre two-room layout. The robot starts at the nurse station, the "
            "patients are at the beds — a near patient and a far patient. As we go left to right, "
            "crowding rises: two people, then four, then seven, all crossing the robot's path. Pay "
            "attention to the far-patient route, because it has to cross this busy corridor — and "
            "that's exactly where the two planners start to behave differently.")

# ============================================================== SLIDE 11 — HOW EACH SOLVES IT
s = add_slide(); title_bar(s, "How Each Planner Handles the Crowd", kicker="Walkthrough")
add_image_fit(s, os.path.join(FIG, "world_high_crowding.png"), Inches(7.0), Inches(1.45),
              Inches(6.0), Inches(5.2))
caption(s, "High crowding — 7 humans crossing the far-patient route", Inches(7.0), Inches(6.7), Inches(6.0))
bullets(s, [
    ("A* hugs the grid — short and predictable…", 0, True),
    ("…but stalls when humans fill its single optimal corridor", 1, False, RED),
    ("RRT* samples freely — finds an alternative path around the crowd…", 0, True),
    ("…and keeps moving toward the patient", 1, False, TEAL),
    ("Guided replay of recorded results — no live run, no network dependency", 0, False, GREY),
], left=Inches(0.55), top=Inches(1.7), width=Inches(6.3), size=19, gap=14)
set_notes(s, "Here's that same crowded scenario. A-star takes the short, direct corridor — optimal on "
            "paper — but when people fill that corridor, it has nowhere to go and stalls. RRT-star, "
            "because it samples the whole space, finds a path around the crowd and keeps moving "
            "toward the patient. So before we even look at numbers, you can see why the stochastic "
            "planner survives the crowd and the deterministic one gets stuck. To be transparent — "
            "this is a guided replay of my recorded results, so nothing here depends on the network "
            "or live hardware.")

# ============================================================== SLIDE 12 — EVAL SETUP
s = add_slide(); title_bar(s, "Evaluation Setup", kicker="Ch 5")
bullets(s, [
    ("Two goals reported separately: near (short route) vs far (long route) — not averaged", 0, True),
    ("they are genuinely different problems", 1, False, GREY),
    ("Metrics: response time · path length · replans · compute time · success rate", 0, False),
    ("All judged against the 60-second deadline", 0, True, TEAL),
    ("10 trials × 2 planners × 3 scenarios", 0, False),
], top=Inches(1.8), size=21, gap=14)
set_notes(s, "On to the evaluation. One important choice: I report the near patient and the far "
            "patient separately instead of averaging them, because they're genuinely different "
            "problems. I measure response time, path length, number of replans, compute time, and "
            "success rate — all against the sixty-second deadline, across ten trials, both planners, "
            "all three crowding levels.")

# ============================================================== SLIDE 13 — RESULTS TABLE
s = add_slide(); title_bar(s, "Results — the Long Route Is Where It Matters", kicker="Ch 5")
sub = s.shapes.add_textbox(Inches(0.7), Inches(1.4), SW - Inches(1.4), Inches(0.5))
tf = sub.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
r = p.add_run(); r.text = "Near patient: planners equivalent — both ~23 s, 100% success everywhere.   "
r.font.size = Pt(18); r.font.color.rgb = DARK
r = p.add_run(); r.text = "Far patient:"
r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = NAVY

rows = [
    ("Scenario", "A* / RRT*  Response (s)", "Success (%)", "Compute (ms)"),
    ("Low",      "51.5 / 41.4", "100 / 100", "9 / 41"),
    ("Moderate", "56.2 / 42.2", "80 / 100",  "13 / 28"),
    ("High",     "77.0 / 54.1", "0 / 100",   "18 / 83"),
]
nrows, ncols = len(rows), 4
tbl_w = Inches(11.0); tbl_h = Inches(3.4)
tx = (SW - tbl_w) // 2; ty = Inches(2.2)
gtbl = s.shapes.add_table(nrows, ncols, tx, ty, tbl_w, tbl_h).table
gtbl.columns[0].width = Inches(2.6)
gtbl.columns[1].width = Inches(3.4)
gtbl.columns[2].width = Inches(2.6)
gtbl.columns[3].width = Inches(2.4)
for ci in range(ncols):
    for ri in range(nrows):
        cell = gtbl.cell(ri, ci)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run(); run.text = rows[ri][ci]
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            run.font.color.rgb = WHITE; run.font.bold = True; run.font.size = Pt(17)
        else:
            highlight = rows[ri][0] == "High"
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF6,0xE2,0xDF) if highlight else (LIGHT if ri % 2 else WHITE)
            run.font.size = Pt(18)
            run.font.bold = highlight
            run.font.color.rgb = RED if highlight else DARK
set_notes(s, "Here are the results. For the near patient, the two planners are basically identical — "
            "both around twenty-three seconds, both succeeding every time. There's no story there. "
            "The story is the far patient. Look at the bottom row, high crowding: A-star takes "
            "seventy-seven seconds and succeeds zero percent of the time — it never beats the "
            "deadline. RRT-star does it in fifty-four seconds and succeeds every single trial.")

# ============================================================== SLIDE 14 — TRADE-OFF + FIGS
s = add_slide(); title_bar(s, "Success Under Deadline & the Trade-off", kicker="Ch 5")
bullets(s, [
    ("A* far-route success fell 100% → 80% → 0% as crowding rose", 0, True, RED),
    ("RRT* met the deadline in every trial (100%)", 0, True, TEAL),
    ("High crowding, A* 0/5 vs RRT* 5/5  →  Fisher's exact p ≈ 0.008 (significant)", 0, False),
    ("Trade-off: RRT* 10–23 s faster, but ~10× more compute & far more variable (83 vs 18 ms)", 0, True),
], left=Inches(0.55), top=Inches(1.55), width=Inches(5.7), size=18, gap=14)
add_image_fit(s, os.path.join(FIG, "success_rate_comparison.png"), Inches(6.5), Inches(1.5),
              Inches(3.3), Inches(5.0))
add_image_fit(s, os.path.join(FIG, "compute_time_comparison.png"), Inches(9.9), Inches(1.5),
              Inches(3.3), Inches(5.0))
set_notes(s, "Putting success rate front and centre: A-star degrades from a hundred percent, to "
            "eighty, to zero as the corridor fills up. RRT-star holds at a hundred throughout. That "
            "difference is statistically significant — a Fisher's exact test gives p around 0.008. "
            "But it's not free. RRT-star is ten to twenty-three seconds faster, yet it costs roughly "
            "ten times more computation and is far less predictable — eighty-three milliseconds "
            "versus eighteen. So it's a genuine reliability-versus-cost trade-off, not a free win.")

# ============================================================== SLIDE 15 — OUTCOMES → OBJECTIVES
s = add_slide(); title_bar(s, "Outcomes Mapped to Objectives", kicker="Ch 5")
pairs = [
    ("Objective", "Outcome"),
    ("1.  Compare A* vs RRT*", "Both built & compared on all metrics"),
    ("2.  Integrate IoT + navigation", "End-to-end: alert → plan → drive → log"),
    ("3.  Evaluate on four metrics", "Per-route, per-crowding, with significance test"),
]
tbl_w = Inches(11.4); tbl_h = Inches(3.6)
tx = (SW - tbl_w) // 2; ty = Inches(2.1)
otbl = s.shapes.add_table(len(pairs), 2, tx, ty, tbl_w, tbl_h).table
otbl.columns[0].width = Inches(4.8); otbl.columns[1].width = Inches(6.6)
for ri, (a, b) in enumerate(pairs):
    for ci, val in enumerate((a, b)):
        cell = otbl.cell(ri, ci); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = cell.text_frame.paragraphs[0]
        run = para.add_run()
        if ri > 0 and ci == 1:
            run.text = "✅  " + val
        else:
            run.text = val
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            run.font.color.rgb = WHITE; run.font.bold = True; run.font.size = Pt(18)
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if ri % 2 else WHITE
            run.font.size = Pt(18); run.font.color.rgb = DARK
            if ci == 0: run.font.bold = True
set_notes(s, "Coming back to my three objectives. One — both planners were implemented and compared "
            "on every metric. Two — the IoT detection and navigation work end to end, from alert to "
            "plan to drive to log. Three — they were evaluated rigorously, broken down by route and "
            "crowding level, with a significance test. All three objectives met.")

# ============================================================== SLIDE 16 — CONCLUSION (narrative arc)
s = add_slide(); title_bar(s, "Conclusion", kicker="Ch 5")
para1 = (
    "This project set out to close two gaps identified in the problem statement: the absence of a "
    "system linking IoT emergency detection to robot navigation, and the common assumption of a "
    "static world in what is really a crowded, dynamic hospital. The aim is met — a complete "
    "sensing-to-navigation pipeline was designed, implemented and evaluated under realistic, "
    "increasingly crowded conditions, with stochastic path planning at its core."
)
para2 = (
    "The results speak to the heart of the problem: in an emergency, a late response is a failed "
    "one. On the critical far-patient route under heavy crowding, the deterministic planner "
    "consistently missed the 60-second clinical deadline — a failed response in every trial — while "
    "the stochastic planner met it every time, a statistically significant difference. The "
    "contribution is therefore not that one planner is universally better, but that stochastic "
    "planning is what keeps emergency response viable precisely when the environment is most "
    "dynamic and the deadline matters most — at a higher but justifiable computational cost."
)
cbox = s.shapes.add_textbox(Inches(0.8), Inches(1.7), SW - Inches(1.6), SH - Inches(2.2))
ctf = cbox.text_frame; ctf.word_wrap = True
for i, para in enumerate((para1, para2)):
    p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
    p.alignment = PP_ALIGN.JUSTIFY
    p.space_after = Pt(18)
    p.line_spacing = 1.18
    r = p.add_run(); r.text = para
    r.font.size = Pt(20); r.font.color.rgb = DARK
set_notes(s, "To conclude, I want to come back to where we started. The problem statement raised two "
            "gaps: there was no system linking IoT emergency detection to robot navigation, and most "
            "navigation work assumes a static world when a hospital is anything but. This project "
            "closes both — I designed, implemented and evaluated a complete sensing-to-navigation "
            "pipeline, with stochastic planning at its core, exactly as the aim set out. And the "
            "results speak to the core of the problem: in an emergency, a late response is a failed "
            "one. On the critical far-patient route under heavy crowding, A-star consistently missed "
            "the sixty-second deadline — a failed response in every single trial — while RRT-star met "
            "it every time, a statistically significant difference. So the contribution isn't that one "
            "planner is universally better; it's that stochastic planning is what keeps emergency "
            "response viable precisely when the environment is most dynamic and the deadline matters "
            "most — at a higher, but justifiable, computational cost. Thank you — I'm happy to take "
            "your questions.")

# ============================================================== SLIDE 17 — Q&A
s = add_slide()
bg = s.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background(); bg.shadow.inherit = False
strip = s.shapes.add_shape(1, Inches(0), Inches(4.05), SW, Inches(0.07))
strip.fill.solid(); strip.fill.fore_color.rgb = TEAL; strip.line.fill.background(); strip.shadow.inherit = False
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.9), SW - Inches(1.8), Inches(1.6))
tf = tb.text_frame
p = tf.paragraphs[0]; r = p.add_run(); r.text = "Thank You"
r.font.size = Pt(48); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph(); r = p2.add_run(); r.text = "Questions & discussion"
r.font.size = Pt(24); r.font.color.rgb = RGBColor(0x9F,0xB6,0xBC)
set_notes(s, "Thank you. I'm happy to take any questions.")

out = os.path.join(ROOT, "defense_ch4_5_slides.pptx")
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
